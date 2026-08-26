# -*- coding: utf-8 -*-
"""
builder.py — собирает готовую презентацию из шаблона.

Поддерживает ДВА типа шаблонов автоматически:

1) "Разметочный" режим — template.pptx с 4-6 слайдами-образцами,
   где текст помечен метками {{TITLE}}, {{TEXT}} и т.д., а картинка —
   фигурой с именем IMAGE_PLACEHOLDER. Первый слайд-образец всегда роль
   "title", последний — всегда "final", всё, что между ними, по порядку
   разбирается на роли из MIDDLE_LAYOUTS_PRIORITY (см. ниже) — так шаблон
   может быть и коротким (4 слайда: text_image + bullets), и длинным
   (6 слайдов: + comparison + highlights).

2) "Автоматический" режим — для обычных PowerPoint-тем без ручной разметки.
   Используются встроенные в тему слайд-макеты (Title Slide,
   Title and Content, Comparison, Picture with Caption и т.д.).
"""
import io
import copy
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PPT
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# Порядок, в котором заполняются "средние" слайды шаблона (между title и
# final). Если в шаблоне 4 слайда — используются первые 2 из этого списка
# (text_image, bullets), если 6 — все 4.
MIDDLE_LAYOUTS_PRIORITY = ["text_image", "bullets", "comparison", "highlights"]


# ==================================================================
#  РЕЖИМ 1: разметочный (метки {{...}} + IMAGE_PLACEHOLDER)
# ==================================================================

def _clone_slide(prs, source_slide):
    # Новый слайд создаём на ТОМ ЖЕ макете, что и слайд-образец — так
    # сохраняется декор, нарисованный на самом макете (фон, рамки, узоры).
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    for shape in source_slide.shapes:
        el = copy.deepcopy(shape._element)
        # Некоторые декоративные фигуры залиты картинкой (не сплошным
        # цветом) — ссылка на эту картинку (blip) действительна только
        # для исходного слайда. Пересоздаём такую связь уже для нового
        # слайда, иначе картинка-заливка пропадёт (будет пустое место).
        for blip in el.iter(qn("a:blip")):
            old_rid = blip.get(qn("r:embed"))
            if not old_rid:
                continue
            try:
                image_part = source_slide.part.related_part(old_rid)
            except KeyError:
                continue
            new_rid = new_slide.part.relate_to(image_part, RT.IMAGE)
            blip.set(qn("r:embed"), new_rid)
        new_slide.shapes._spTree.insert_element_before(el, "p:extLst")
    return new_slide


def _fill_slide_markers(slide, data: dict, image_bytes):
    for shape in list(slide.shapes):
        if shape.name == "IMAGE_PLACEHOLDER":
            if image_bytes:
                left, top, w, h = shape.left, shape.top, shape.width, shape.height
                shape._element.getparent().remove(shape._element)
                slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, w, h)
            else:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for r in p.runs:
                            r.text = ""
            continue
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for key, value in data.items():
                        if isinstance(value, str):
                            run.text = run.text.replace("{{" + key + "}}", value)
            # если ИИ-текст длиннее, чем рассчитывал автор шаблона — ужимаем
            # шрифт под рамку, а не даём ему наезжать на соседние блоки
            try:
                shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass


def _has_markers(prs) -> bool:
    for slide in list(prs.slides)[:6]:
        for shape in slide.shapes:
            if shape.has_text_frame and "{{" in shape.text_frame.text:
                return True
    return False


def _layout_plan_from_prs(prs) -> dict:
    """
    Определяет, какая роль соответствует какому слайду-образцу в файле:
    первый — всегда title, последний — всегда final, остальное — по
    порядку из MIDDLE_LAYOUTS_PRIORITY. Работает для 4, 5, 6... слайдов.
    """
    n = len(list(prs.slides))
    plan = {"title": 0}
    middle_count = max(0, n - 2)
    for i, role in enumerate(MIDDLE_LAYOUTS_PRIORITY[:middle_count]):
        plan[role] = i + 1
    if n >= 2:
        plan["final"] = n - 1
    return plan


def get_available_layouts(template_path: str) -> list:
    """
    Список ролей слайдов, которые реально поддерживает этот шаблон.
    Используется перед обращением к нейросети — чтобы попросить структуру
    только из тех типов слайдов, что есть в файле (иначе, например,
    4-слайдовый шаблон получит запрос на несуществующий 5-й тип и упадёт).
    """
    prs = Presentation(template_path)
    if _has_markers(prs):
        plan = _layout_plan_from_prs(prs)
        order = ["title"] + MIDDLE_LAYOUTS_PRIORITY + ["final"]
        return [role for role in order if role in plan]
    else:
        return ["title", "text_image", "bullets", "final"]


def _build_marker_mode(prs, structure, images, layout_plan: dict) -> bytes:
    template_slides = list(prs.slides)
    fallback_idx = layout_plan.get("bullets", 1)

    for i, slide_data in enumerate(structure["slides"]):
        role = slide_data.get("layout", "bullets")
        layout_idx = layout_plan.get(role, fallback_idx)
        new_slide = _clone_slide(prs, template_slides[layout_idx])
        _fill_slide_markers(new_slide, slide_data, images.get(i))

    xml_slides = prs.slides._sldIdLst
    sld_ids = list(xml_slides)
    for sld_id in sld_ids[:len(template_slides)]:
        prs.part.drop_rel(sld_id.rId)
        xml_slides.remove(sld_id)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ==================================================================
#  РЕЖИМ 2: автоматический (родные слайд-макеты темы)
# ==================================================================

_HOUSEKEEPING = {PPT.DATE, PPT.FOOTER, PPT.SLIDE_NUMBER, PPT.HEADER}


def _pick_layout(prs, groups: list, avoid: set = frozenset()):
    candidates = []
    for layout in prs.slide_layouts:
        types = {ph.placeholder_format.type for ph in layout.placeholders}
        if types & avoid:
            continue
        matched = sum(1 for g in groups if types & g)
        if matched == 0:
            continue
        content_size = len(types - _HOUSEKEEPING)
        candidates.append((matched, content_size, layout))

    if not candidates:
        return prs.slide_layouts[0]

    candidates.sort(key=lambda t: (-t[0], t[1]))
    return candidates[0][2]


def _find_ph(slide, types: set):
    for ph in slide.placeholders:
        if ph.placeholder_format.type in types:
            return ph
    return None


def _find_all_ph(slide, types: set) -> list:
    return [ph for ph in slide.placeholders if ph.placeholder_format.type in types]


def _build_native_slide(prs, slide_data: dict, image_bytes):
    kind = slide_data.get("layout", "bullets")

    if kind == "title":
        layout = _pick_layout(prs, [{PPT.CENTER_TITLE, PPT.TITLE}, {PPT.SUBTITLE}],
                               avoid={PPT.OBJECT, PPT.BODY, PPT.PICTURE})
    elif kind == "text_image":
        layout = _pick_layout(prs, [{PPT.TITLE, PPT.CENTER_TITLE}, {PPT.PICTURE}])
    elif kind == "final":
        layout = _pick_layout(prs, [{PPT.TITLE, PPT.CENTER_TITLE}],
                               avoid={PPT.OBJECT, PPT.BODY, PPT.PICTURE})
    elif kind in ("comparison", "highlights"):
        # ищем макет с ДВУМЯ текстовыми полями (например, "Comparison" в
        # стандартной теме) — если такого нет, сработает общий fallback ниже
        layout = _pick_layout(prs, [{PPT.TITLE, PPT.CENTER_TITLE}, {PPT.BODY, PPT.OBJECT}],
                               avoid={PPT.PICTURE})
    else:  # bullets и всё неизвестное
        layout = _pick_layout(prs, [{PPT.TITLE, PPT.CENTER_TITLE}, {PPT.BODY, PPT.OBJECT}],
                               avoid={PPT.PICTURE})

    slide = prs.slides.add_slide(layout)

    title_ph = _find_ph(slide, {PPT.TITLE, PPT.CENTER_TITLE})
    if title_ph and slide_data.get("TITLE"):
        title_ph.text_frame.text = slide_data["TITLE"]

    sub_ph = _find_ph(slide, {PPT.SUBTITLE})
    if sub_ph and slide_data.get("SUBTITLE"):
        sub_ph.text_frame.text = slide_data["SUBTITLE"]

    if kind in ("comparison", "highlights"):
        body_phs = _find_all_ph(slide, {PPT.BODY, PPT.OBJECT})
        pair = [slide_data.get("TEXT_1"), slide_data.get("TEXT_2")] if kind == "comparison" \
            else [slide_data.get("HIGHLIGHT_1"), slide_data.get("HIGHLIGHT_2")]
        for ph, text in zip(body_phs, pair):
            if text:
                ph.text_frame.text = text
    else:
        body_ph = _find_ph(slide, {PPT.BODY, PPT.OBJECT})
        bullets = [v for k, v in sorted(slide_data.items()) if k.startswith("BULLET_") and v]
        if body_ph and bullets:
            tf = body_ph.text_frame
            tf.clear()
            for i, b in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = b
        elif body_ph and slide_data.get("TEXT"):
            body_ph.text_frame.text = slide_data["TEXT"]

    pic_ph = _find_ph(slide, {PPT.PICTURE})
    if image_bytes:
        if pic_ph:
            pic_ph.insert_picture(io.BytesIO(image_bytes))
        else:
            slide.shapes.add_picture(
                io.BytesIO(image_bytes),
                Emu(int(prs.slide_width * 0.55)),
                Emu(int(prs.slide_height * 0.15)),
                height=Emu(int(prs.slide_height * 0.7)),
            )


def _build_native_mode(prs, structure, images) -> bytes:
    original_count = len(list(prs.slides))

    for i, slide_data in enumerate(structure["slides"]):
        _build_native_slide(prs, slide_data, images.get(i))

    xml_slides = prs.slides._sldIdLst
    sld_ids = list(xml_slides)
    for sld_id in sld_ids[:original_count]:
        prs.part.drop_rel(sld_id.rId)
        xml_slides.remove(sld_id)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ==================================================================
#  ТОЧКА ВХОДА
# ==================================================================

def build_pptx(template_path: str, structure: dict, images: dict) -> bytes:
    prs = Presentation(template_path)

    if _has_markers(prs):
        plan = _layout_plan_from_prs(prs)
        return _build_marker_mode(prs, structure, images, plan)
    else:
        return _build_native_mode(prs, structure, images)
